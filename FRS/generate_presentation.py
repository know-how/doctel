#!/usr/bin/env python3
"""
Generate DocTel FRS PowerPoint Presentation
Creates a comprehensive presentation with slides covering:
- System Overview
- Architecture
- Functional Requirements
- Data Flows
- User Workflows
- Technical Details
- Security & Performance
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(91, 136, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"May 10, 2026"
    date_frame.paragraphs[0].font.size = Pt(14)
    date_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 231, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.level = 0
        p.space_before = Pt(12)
    
    return slide

def create_presentation():
    """Create complete DocTel FRS presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, "DocTel (DocIntel)", "Document Analysis & AI-Powered Insights System\nFunctional Requirement Specification")
    
    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "🎯 Purpose: Enable document intelligence through AI-powered analysis",
        "📤 Upload: Support PDF, DOCX, TXT formats with metadata tagging",
        "🔍 Analysis: Auto-generate summaries, entities, sentiment, topics",
        "💡 Prompts: Generate 3-5 context-specific questions per document",
        "❓ Q&A: RAG-powered answers with source citations",
        "🔐 Security: Role-based access control, encrypted storage"
    ])
    
    # Slide 3: Key Features
    add_content_slide(prs, "Core Features", [
        "✓ Multi-Format Document Support (PDF, DOCX, TXT)",
        "✓ Automatic Executive Summaries (up to 10 sentences)",
        "✓ Detailed Analysis with Bullet Points & Key Insights",
        "✓ Entity Extraction (People, Places, Dates, Concepts)",
        "✓ Sentiment Analysis (Positive, Neutral, Negative, Urgent)",
        "✓ AI-Generated Suggested Questions (3-5 per document)",
        "✓ Chat Interface with RAG-Powered Answers & Citations"
    ])
    
    # Slide 4: System Architecture
    add_content_slide(prs, "System Architecture", [
        "🌐 Frontend: React + Vite (Web) | React Native (Mobile)",
        "⚙️  Backend: FastAPI (Python) with async/await processing",
        "🗄️  Database: PostgreSQL + Chroma Vector Store",
        "🤖 LLM Layer: Ollama (local) + Gemini API (cloud fallback)",
        "🔐 Auth: Email OTP + Active Directory integration",
        "📊 Model Router: Intelligent fallback chain selection",
        "⚡ Processing: Async ingestion worker for concurrent jobs"
    ])
    
    # Slide 5: Document Ingestion Pipeline
    add_content_slide(prs, "Document Ingestion Pipeline", [
        "1️⃣  User uploads document (PDF/DOCX/TXT)",
        "2️⃣  System validates format and generates hash (SHA256)",
        "3️⃣  Extract text using format-specific parser",
        "4️⃣  Split into chunks (500-1000 tokens with overlap)",
        "5️⃣  Generate embeddings for each chunk",
        "6️⃣  Store embeddings in Chroma vector database",
        "7️⃣  Trigger analysis (summaries, entities, sentiment, prompts)",
        "8️⃣  Mark document as 'analysis_ready = true'"
    ])
    
    # Slide 6: Analysis Capabilities
    add_content_slide(prs, "Document Analysis Capabilities", [
        "📝 Executive Summary: Concise 5-10 sentence overview",
        "📋 Detailed Summary: Structured sections with bullet points",
        "👥 Entity Extraction: People, organizations, dates, locations",
        "😊 Sentiment Analysis: Positive/Neutral/Negative/Urgent classification",
        "🏷️  Topic Extraction: Main themes with relevance scores",
        "✅ Action Items: Extracted tasks with owners and deadlines",
        "🎯 Decisions: Approved actions and policy changes"
    ])
    
    # Slide 7: Prompt Generation
    add_content_slide(prs, "Intelligent Prompt Generation", [
        "🔧 Automatic Generation: 3-5 context-specific questions",
        "📚 Document Type Aware: Different questions for policies vs reports",
        "🎯 Domain Specific: ZETDC terminology (SCADA, transformers, feeders)",
        "💡 Examples:",
        "   • Policy: 'What are the compliance requirements?'",
        "   • Report: 'What are the key metrics?'",
        "   • Minutes: 'What action items were assigned?'",
        "✨ Users click questions or type custom queries"
    ])
    
    # Slide 8: RAG-Powered Q&A
    add_content_slide(prs, "RAG Question Answering with Citations", [
        "1️⃣  User asks question about document",
        "2️⃣  System converts question to embedding vector",
        "3️⃣  Vector search in Chroma DB retrieves top-K chunks",
        "4️⃣  Chunks ranked by semantic relevance",
        "5️⃣  Context assembled from relevant chunks",
        "6️⃣  LLM generates answer using context",
        "7️⃣  Answer includes inline citations: [Doc: filename, chunk N]",
        "8️⃣  Citations clickable to view source text"
    ])
    
    # Slide 9: User Workflow - Upload
    add_content_slide(prs, "User Workflow: Upload & Analyze", [
        "📤 Step 1: User clicks 'Upload Document' (or drag-drop)",
        "📋 Step 2: Select file, project, document type, add tags",
        "⏳ Step 3: System processes (progress bar shows: Extracting → Embedding → Summarizing)",
        "📊 Step 4: View analysis dashboard with:",
        "   • Executive summary",
        "   • Key insights (entities, sentiment, themes, actions)",
        "   • 3-5 suggested questions",
        "✅ Step 5: Document ready for Q&A"
    ])
    
    # Slide 10: User Workflow - Ask Questions
    add_content_slide(prs, "User Workflow: Ask Questions", [
        "📖 Option A - Use Suggested Prompt:",
        "   • Click one of 3-5 pre-generated questions",
        "   • System auto-submits and shows answer",
        "",
        "✏️  Option B - Ask Custom Question:",
        "   • Type question in 'Ask about this document...' box",
        "   • Press Enter",
        "   • System searches document and generates answer with citations",
        "💬 Step 4: Chat history maintained for session"
    ])
    
    # Slide 11: Project-Level Analysis
    add_content_slide(prs, "Project-Level Features", [
        "📁 Group Related Documents: Create projects for themes/periods",
        "📤 Bulk Upload: Upload 3+ related documents together",
        "🔍 Cross-Document Search: Ask questions about entire project",
        "📊 Comparison Analysis: Compare metrics/findings across documents",
        "📑 Multi-Document Reports: Generate synthesis reports",
        "🔗 Document Relationships: Track connections between documents",
        "👥 Team Collaboration: Share projects with team members"
    ])
    
    # Slide 12: Database Schema
    add_content_slide(prs, "Data Model Overview", [
        "👤 Users: Username, email, role (admin/analyst/viewer)",
        "📁 Projects: Contains related documents, has members",
        "📄 Documents: File metadata, ingestion status, analysis results",
        "📊 DocAnalysis: Summaries, entities, sentiment, topics, actions",
        "🧩 Chunks: Text pieces with embeddings for vector search",
        "🎯 SuggestedPrompts: 3-5 questions generated per document",
        "💬 Sessions & Messages: Chat history with citations"
    ])
    
    # Slide 13: Security Model
    add_content_slide(prs, "Security & Access Control", [
        "🔐 Authentication: Email OTP (primary) or Active Directory",
        "👮 Role-Based Access (RBAC):",
        "   • Admin: Full system control",
        "   • Analyst: Create/upload documents",
        "   • Viewer: Read-only access",
        "🔒 Data Protection: TLS encryption in transit, configurable at-rest",
        "📋 Audit Logging: All document access logged with timestamp/user",
        "🔑 API Security: Token-based auth (24-hour expiry)"
    ])
    
    # Slide 14: Performance Metrics
    add_content_slide(prs, "Performance Requirements", [
        "⚡ Document Upload: 2-5 seconds (10-page PDF)",
        "⚡ Text Extraction: 3-10 seconds",
        "⚡ Embedding Generation: 10-20 seconds per 10 chunks",
        "⚡ Complete Analysis: 30-60 seconds for 10-page document",
        "⚡ RAG Answer: 3-10 seconds (vector search + LLM)",
        "🚀 Concurrency: Support 100+ concurrent users",
        "🔄 Processing: 5 documents ingesting simultaneously"
    ])
    
    # Slide 15: LLM Integration
    add_content_slide(prs, "LLM & Model Management", [
        "🤖 Primary: Ollama (Local models - privacy-first)",
        "☁️  Fallback: Gemini API (when local unavailable)",
        "🔄 Model Router: Intelligent selection with fallback chain",
        "📦 Supported Models: Llama2, Mistral, Neural-Chat, etc.",
        "⚙️  Configuration: Per-task model selection + temperature",
        "💾 Model Cache: Local model inventory with status",
        "🚀 Auto-Fallback: If model unavailable, try next in chain"
    ])
    
    # Slide 16: Error Handling
    add_content_slide(prs, "Error Handling & Recovery", [
        "❌ Unsupported Format: Reject with 422 status + helpful message",
        "❌ Embedding Failed: Auto-retry up to 3 times, then queue",
        "❌ Model Unavailable: Fallback to next model in chain",
        "❌ DB Connection Lost: Automatic retry with exponential backoff",
        "❌ No Relevant Chunks: Helpful message + suggest rephrasing",
        "❌ Document Not Found: 404 with list of available documents",
        "✅ All errors logged for debugging and monitoring"
    ])
    
    # Slide 17: Technology Stack
    add_content_slide(prs, "Technology Stack", [
        "🐍 Backend: Python 3.10+ with FastAPI framework",
        "🗄️  Database: PostgreSQL 14+ with SQLAlchemy ORM",
        "🔍 Vector Store: Chroma for embeddings and semantic search",
        "🎨 Frontend: React 18 + Vite + TypeScript",
        "📱 Mobile: React Native with Expo",
        "🤖 LLM: Ollama + Gemini API",
        "🔐 Auth: Email (SMTP) + LDAP/AD integration"
    ])
    
    # Slide 18: Deployment Architecture
    add_content_slide(prs, "Deployment Architecture", [
        "🔄 Load Balancer: Nginx distributing traffic",
        "🏢 App Servers: Multiple FastAPI instances (horizontal scaling)",
        "💾 Database: PostgreSQL with connection pooling",
        "🤖 LLM Services: Ollama containers + Gemini cloud API",
        "🧬 Vector DB: Chroma instance shared by app servers",
        "⚙️  Background: Ingestion worker processes documents async",
        "📊 Monitoring: Logging, metrics, alerts (configurable)"
    ])
    
    # Slide 19: API Endpoints
    add_content_slide(prs, "Key API Endpoints", [
        "📤 POST /api/documents/upload - Upload new document",
        "📖 GET /api/documents/{doc_id} - Get document details",
        "❓ POST /api/documents/{doc_id}/ask - Ask about document",
        "🔍 POST /api/projects/{proj_id}/ask - Cross-document query",
        "📁 GET /api/projects - List user's projects",
        "📊 GET /api/models/available - List available models",
        "💬 GET /api/sessions/{sess_id}/messages - Chat history"
    ])
    
    # Slide 20: Data Flow Diagram
    add_content_slide(prs, "Document Ingestion Data Flow", [
        "[Upload] → [Validate] → [Store DB] → [Extract Text]",
        "  ↓         ↓          ↓          ↓",
        "[Split into Chunks] → [Generate Embeddings] → [Chroma]",
        "  ↓",
        "[Generate Analysis]",
        "  ├─ Executive Summary",
        "  ├─ Entity Extraction",
        "  ├─ Sentiment Analysis",
        "  ├─ Topic Extraction",
        "  └─ Generate Prompts",
        "  ↓",
        "[Ready for Q&A]"
    ])
    
    # Slide 21: Q&A Data Flow
    add_content_slide(prs, "Question Answering Data Flow", [
        "[User Question] → [Embed Question]",
        "  ↓",
        "[Vector Search in Chroma] → [Retrieve Top Chunks]",
        "  ↓",
        "[Assemble Context] → [Select LLM Model]",
        "  ↓",
        "[Generate Answer] ← [Ollama/Gemini]",
        "  ↓",
        "[Extract Citations] → [Store in Chat History]",
        "  ↓",
        "[Return to User with Citations]"
    ])
    
    # Slide 22: Configuration
    add_content_slide(prs, "System Configuration", [
        "⚙️  Embedding Model: nomic-embed-text (384 dimensions)",
        "⚙️  Chunk Size: 1000 tokens with 200-token overlap",
        "⚙️  Vector Search K: 6 most relevant chunks",
        "⚙️  LLM Temperature: 0.7 (balanced creativity/accuracy)",
        "⚙️  Token Expiry: 24 hours",
        "⚙️  OTP Validity: 15 minutes",
        "⚙️  Max File Size: 100 MB"
    ])
    
    # Slide 23: Use Cases
    add_content_slide(prs, "Real-World Use Cases", [
        "📋 Policy Review: Upload company policies, ask specific questions",
        "📊 Report Analysis: Summarize quarterly reports, compare metrics",
        "📝 Meeting Minutes: Extract action items and decisions",
        "🔍 Compliance Audit: Search all policies for specific requirements",
        "🎓 Knowledge Base: Team learns from document repository",
        "📑 Cross-Reference: Find related documents automatically",
        "🏗️  Infrastructure Planning: Analyze technical specifications"
    ])
    
    # Slide 24: Future Enhancements
    add_content_slide(prs, "Phase 2 & 3 Roadmap", [
        "🔮 Multi-Document Reports: Auto-generate synthesis documents",
        "🎨 Diagram Generation: Auto-create flowcharts, network diagrams",
        "📋 Policy Generation: AI-assisted policy draft creation",
        "🤖 Fine-Tuned Models: Domain-specific model training",
        "🔗 Integration APIs: Salesforce, SharePoint, Slack connectors",
        "📱 Full Mobile App: React Native offline support",
        "🌐 Global Search: Full-text + semantic search across all docs"
    ])
    
    # Slide 25: Business Impact
    add_content_slide(prs, "Expected Business Impact", [
        "⏱️  70% Reduction in document review time",
        "🔍 90% Improvement in information discoverability",
        "⚡ Faster decision-making with AI insights",
        "✅ Better compliance with audit trails",
        "📚 Knowledge preservation and dissemination",
        "💰 ROI within 6 months through time savings",
        "🏆 Competitive advantage in document intelligence"
    ])
    
    # Slide 26: Implementation Timeline
    add_content_slide(prs, "Implementation Timeline", [
        "📅 Week 1-2: Environment setup & infrastructure",
        "📅 Week 3-4: Core backend services (ingestion, RAG)",
        "📅 Week 5-6: Frontend development (React)",
        "📅 Week 7-8: Mobile app (React Native)",
        "📅 Week 9-10: Testing & optimization",
        "📅 Week 11-12: Deployment & UAT",
        "📅 Week 13+: Production & Phase 2 planning"
    ])
    
    # Slide 27: Risk Mitigation
    add_content_slide(prs, "Risk Mitigation Strategies", [
        "⚠️  Model Unavailability: Fallback chain + error handling",
        "⚠️  Data Breach: Encryption + RBAC + audit logging",
        "⚠️  Performance Issues: Load testing + optimization",
        "⚠️  User Adoption: Training + documentation + support",
        "⚠️  Integration Issues: API testing + staged rollout",
        "⚠️  Scaling Challenges: Horizontal scaling architecture",
        "⚠️  Model Accuracy: Human review + feedback loops"
    ])
    
    # Slide 28: Key Metrics
    add_content_slide(prs, "Success Metrics & KPIs", [
        "📈 Document Processing Rate: >1000 docs/month",
        "📈 User Adoption: >80% of target users within 3 months",
        "📈 Average Response Time: <5 seconds for Q&A",
        "📈 System Uptime: >99.5%",
        "📈 User Satisfaction: >4.5/5 star rating",
        "📈 Accuracy: >90% on citation relevance",
        "📈 ROI: 6-month payback period"
    ])
    
    # Slide 29: Questions & Support
    add_content_slide(prs, "Questions & Support", [
        "📞 Technical Support: IT Help Desk",
        "📧 Feature Requests: Product Team",
        "🐛 Bug Reports: DevOps Team",
        "📚 Documentation: Wiki + Help Center",
        "👥 Training: Hands-on sessions scheduled",
        "🎓 Certification: DocTel Power User program",
        "🔗 Community: Internal Slack channel"
    ])
    
    # Slide 30: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    conclusion_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    conclusion_frame = conclusion_box.text_frame
    conclusion_frame.word_wrap = True
    
    conclusion_text = """DocTel represents a transformational 
    shift in how organizations interact with 
    their document repositories.
    
    By combining modern AI, NLP, and enterprise 
    architecture, we deliver a system that is:
    
    Intelligent • Scalable • Secure • User-Friendly
    
    Ready to empower better decisions through 
    document intelligence."""
    
    conclusion_frame.text = conclusion_text
    for paragraph in conclusion_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(200, 200, 200)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = r"c:\Users\ze9167523\IdeaProjects\doctel\FRS\DocTel_FRS_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_presentation()
